#!/usr/bin/env python3
"""
Docling Document Extractor for TechWriterReview
================================================
Version: reads from version.json (module v1.1)
Date: 2026-01-27

Advanced document parsing using Docling (IBM's open-source document parser).
Fully optimized for AIR-GAPPED NETWORKS with NO internet connectivity.

Features:
- Unified extraction for PDF, DOCX, PPTX, XLSX, HTML
- Advanced table structure recognition (AI-powered)
- Reading order preservation
- Memory-optimized (images excluded by default)
- Complete offline operation when properly configured

RUNTIME vs INSTALLATION:
- INSTALLATION: Requires internet (use setup_docling.bat)
- RUNTIME: 100% offline - environment variables block network access

To allow network access (e.g., for model updates), set:
    DOCLING_ALLOW_NETWORK=1

Author: Nick / SAIC Systems Engineering
"""

import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple, Union
from dataclasses import dataclass, field
from enum import Enum
import json
import logging
import time

__version__ = "1.1.0"

# ============================================================================
# AIR-GAP CONFIGURATION - Set these BEFORE importing docling
# ============================================================================
# These settings ensure Docling operates 100% offline at RUNTIME.
# They do NOT affect installation - setup_docling.bat uses internet to download.
#
# To disable offline mode (allow network access), set environment variable:
#   DOCLING_ALLOW_NETWORK=1
# ============================================================================

def _configure_offline_environment():
    """
    Configure environment for offline operation.
    
    Only applies if DOCLING_ALLOW_NETWORK is not set.
    This blocks Hugging Face libraries from making network requests.
    """
    # Check if user wants to allow network access
    if os.environ.get('DOCLING_ALLOW_NETWORK', '').lower() in ('1', 'true', 'yes'):
        return  # Skip offline configuration
    
    # Force offline mode - prevents ANY network access during runtime
    os.environ.setdefault('HF_HUB_OFFLINE', '1')  # Hugging Face offline
    os.environ.setdefault('TRANSFORMERS_OFFLINE', '1')  # Transformers offline
    os.environ.setdefault('HF_DATASETS_OFFLINE', '1')  # Datasets offline

    # Disable telemetry/analytics
    os.environ.setdefault('HF_HUB_DISABLE_TELEMETRY', '1')
    os.environ.setdefault('DO_NOT_TRACK', '1')
    os.environ.setdefault('ANONYMIZED_TELEMETRY', 'false')

# Apply offline configuration at import time (runtime only)
_configure_offline_environment()

# ============================================================================
# LOGGING
# ============================================================================

try:
    from config_logging import get_logger
    _logger = get_logger('docling_extractor')
except ImportError:
    _logger = logging.getLogger('docling_extractor')
    if not _logger.handlers:
        handler = logging.StreamHandler()
        handler.setFormatter(logging.Formatter('[%(name)s] %(levelname)s: %(message)s'))
        _logger.addHandler(handler)
        _logger.setLevel(logging.INFO)


def _log(message: str, level: str = 'info', **kwargs):
    """Internal logging helper."""
    getattr(_logger, level)(message)


# ============================================================================
# DATA CLASSES
# ============================================================================

class ExtractionBackend(Enum):
    """Available extraction backends."""
    DOCLING = "docling"
    LEGACY = "legacy"
    AUTO = "auto"


@dataclass
class ExtractedTable:
    """Represents an extracted table from a document."""
    table_id: int
    page_number: int
    rows: List[List[str]]
    headers: List[str] = field(default_factory=list)
    caption: str = ""
    confidence: float = 0.0
    
    @property
    def row_count(self) -> int:
        return len(self.rows)
    
    @property
    def col_count(self) -> int:
        return len(self.rows[0]) if self.rows else len(self.headers)
    
    def to_csv_rows(self) -> List[List[str]]:
        """Return table as CSV-compatible rows."""
        if self.headers:
            return [self.headers] + self.rows
        return self.rows
    
    def to_dict(self) -> dict:
        return {
            'table_id': self.table_id,
            'page_number': self.page_number,
            'headers': self.headers,
            'rows': self.rows,
            'caption': self.caption,
            'confidence': self.confidence,
            'row_count': self.row_count,
            'col_count': self.col_count
        }


@dataclass
class ExtractedSection:
    """Represents a document section/heading."""
    level: int
    title: str
    content: str
    page_number: int
    
    def to_dict(self) -> dict:
        return {
            'level': self.level,
            'title': self.title,
            'content': self.content,
            'page_number': self.page_number
        }


@dataclass
class ExtractedParagraph:
    """Represents a text block with location metadata."""
    text: str
    location: str
    page_number: int = 1
    paragraph_type: str = "text"  # text, heading, list_item, table_cell
    confidence: float = 1.0
    
    def to_dict(self) -> dict:
        return {
            'text': self.text,
            'location': self.location,
            'page_number': self.page_number,
            'type': self.paragraph_type,
            'confidence': self.confidence
        }


@dataclass
class DocumentExtractionResult:
    """Complete extraction result from a document."""
    filepath: str
    filename: str
    format: str
    backend_used: str
    
    # Content
    full_text: str = ""
    markdown: str = ""
    sections: List[ExtractedSection] = field(default_factory=list)
    tables: List[ExtractedTable] = field(default_factory=list)
    paragraphs: List[ExtractedParagraph] = field(default_factory=list)
    
    # Metadata
    page_count: int = 0
    word_count: int = 0
    char_count: int = 0
    extraction_time_ms: float = 0.0
    warnings: List[str] = field(default_factory=list)
    
    # Docling-specific metadata
    docling_version: str = ""
    models_used: List[str] = field(default_factory=list)
    offline_mode: bool = True
    
    def to_dict(self) -> dict:
        return {
            'filepath': self.filepath,
            'filename': self.filename,
            'format': self.format,
            'backend_used': self.backend_used,
            'page_count': self.page_count,
            'word_count': self.word_count,
            'char_count': self.char_count,
            'table_count': len(self.tables),
            'section_count': len(self.sections),
            'paragraph_count': len(self.paragraphs),
            'extraction_time_ms': self.extraction_time_ms,
            'warnings': self.warnings,
            'docling_version': self.docling_version,
            'models_used': self.models_used,
            'offline_mode': self.offline_mode
        }
    
    def get_paragraphs_for_roles(self) -> List[Tuple[int, str]]:
        """Return paragraphs in format expected by role extractor."""
        return [(i, p.text) for i, p in enumerate(self.paragraphs)]
    
    def get_text_by_page(self) -> Dict[int, str]:
        """Return text grouped by page number."""
        pages: Dict[int, List[str]] = {}
        for p in self.paragraphs:
            if p.page_number not in pages:
                pages[p.page_number] = []
            pages[p.page_number].append(p.text)
        return {page: '\n'.join(texts) for page, texts in pages.items()}


# ============================================================================
# DOCLING EXTRACTOR - MAIN CLASS
# ============================================================================

class DoclingExtractor:
    """
    Document extraction using Docling library.
    
    Optimized for AIR-GAPPED environments with:
    - Complete offline operation
    - Memory optimization (no image processing)
    - Maximum text/table extraction quality
    """
    
    SUPPORTED_FORMATS = {
        '.pdf', '.docx', '.doc', '.pptx', '.ppt', 
        '.xlsx', '.xls', '.html', '.htm', '.md', '.txt'
    }
    
    # Default configuration optimized for air-gapped + memory efficiency
    DEFAULT_CONFIG = {
        'artifacts_path': None,  # Set from env or parameter
        'enable_ocr': False,     # Disabled by default - enable if needed
        'ocr_engine': 'easyocr', # Preferred OCR if enabled
        'table_mode': 'accurate', # 'fast' or 'accurate'
        'do_picture_classifer': False,  # Disabled - saves memory
        'do_picture_description': False, # Disabled - saves memory
        'generate_page_images': False,   # Disabled - saves memory
        'generate_picture_images': False, # Disabled - saves memory
        'images_scale': 1.0,     # Not used when images disabled
        'enable_remote_services': False,  # CRITICAL: No network access
        'fallback_to_legacy': True,
    }
    
    def __init__(
        self,
        artifacts_path: Optional[str] = None,
        enable_ocr: bool = False,
        ocr_engine: str = 'easyocr',
        table_mode: str = 'accurate',
        fallback_to_legacy: bool = True,
        **kwargs
    ):
        """
        Initialize the Docling extractor for air-gapped operation.
        
        Args:
            artifacts_path: Path to pre-downloaded Docling models.
                           If None, uses DOCLING_ARTIFACTS_PATH env var.
            enable_ocr: Enable OCR for scanned documents (requires models)
            ocr_engine: OCR engine ('easyocr', 'tesseract', 'rapidocr')
            table_mode: Table extraction mode ('fast' or 'accurate')
            fallback_to_legacy: Fall back to pdfplumber/python-docx on errors
        """
        # Resolve artifacts path
        self.artifacts_path = (
            artifacts_path or 
            os.environ.get('DOCLING_ARTIFACTS_PATH') or
            os.environ.get('DOCLING_SERVE_ARTIFACTS_PATH')
        )
        
        self.enable_ocr = enable_ocr
        self.ocr_engine = ocr_engine
        self.table_mode = table_mode
        self.fallback_to_legacy = fallback_to_legacy
        
        # Merge any additional kwargs with defaults
        self.config = {**self.DEFAULT_CONFIG, **kwargs}
        self.config['artifacts_path'] = self.artifacts_path
        self.config['enable_ocr'] = enable_ocr
        
        self._docling_available = False
        self._docling_version = ""
        self._converter = None
        self._pdf_pipeline_options = None
        
        self._init_docling()
    
    def _init_docling(self):
        """Initialize Docling with air-gapped configuration."""
        try:
            # Import Docling components
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions,
                TableFormerMode,
            )
            
            # Get Docling version
            try:
                import docling
                self._docling_version = getattr(docling, '__version__', 'unknown')
            except:
                self._docling_version = 'unknown'
            
            # Configure PDF pipeline with air-gap optimizations
            pipeline_options = PdfPipelineOptions()
            
            # === CRITICAL: Air-gapped configuration ===
            if self.artifacts_path:
                pipeline_options.artifacts_path = self.artifacts_path
                _log(f"Docling using local models: {self.artifacts_path}")
            
            # Disable remote services (no network access)
            if hasattr(pipeline_options, 'enable_remote_services'):
                pipeline_options.enable_remote_services = False
            
            # === Table extraction (enabled - this is a key feature) ===
            pipeline_options.do_table_structure = True
            if hasattr(pipeline_options, 'table_structure_options'):
                if self.table_mode == 'accurate':
                    pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE
                else:
                    pipeline_options.table_structure_options.mode = TableFormerMode.FAST
                # Use text cells from PDF when available (better accuracy)
                pipeline_options.table_structure_options.do_cell_matching = True
            
            # === Memory optimization: Disable image processing ===
            if hasattr(pipeline_options, 'do_picture_classifier'):
                pipeline_options.do_picture_classifier = False
            if hasattr(pipeline_options, 'do_picture_description'):
                pipeline_options.do_picture_description = False
            if hasattr(pipeline_options, 'generate_page_images'):
                pipeline_options.generate_page_images = False
            if hasattr(pipeline_options, 'generate_picture_images'):
                pipeline_options.generate_picture_images = False
            if hasattr(pipeline_options, 'images_scale'):
                pipeline_options.images_scale = 1.0  # Minimum if images were processed
            
            # === OCR configuration (optional) ===
            pipeline_options.do_ocr = self.enable_ocr
            if self.enable_ocr:
                self._configure_ocr(pipeline_options)
            
            # Store pipeline options for reference
            self._pdf_pipeline_options = pipeline_options
            
            # Create the converter
            self._converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
            
            self._docling_available = True
            _log(f"Docling {self._docling_version} initialized (offline mode, "
                 f"table_mode={self.table_mode}, ocr={self.enable_ocr})")
            
        except ImportError as e:
            _log(f"Docling not available: {e}", level='warning')
            self._docling_available = False
        except Exception as e:
            _log(f"Error initializing Docling: {e}", level='error')
            self._docling_available = False
    
    def _configure_ocr(self, pipeline_options):
        """Configure OCR engine for air-gapped operation."""
        try:
            if self.ocr_engine == 'easyocr':
                from docling.datamodel.pipeline_options import EasyOcrOptions
                ocr_options = EasyOcrOptions(lang=['en'])
                # EasyOCR models should be in artifacts_path
                if self.artifacts_path:
                    ocr_options.model_storage_directory = self.artifacts_path
                pipeline_options.ocr_options = ocr_options
                
            elif self.ocr_engine == 'tesseract':
                from docling.datamodel.pipeline_options import TesseractOcrOptions
                pipeline_options.ocr_options = TesseractOcrOptions(lang=['eng'])
                
            elif self.ocr_engine == 'rapidocr':
                from docling.datamodel.pipeline_options import RapidOcrOptions
                ocr_options = RapidOcrOptions()
                if self.artifacts_path:
                    ocr_options.model_path = self.artifacts_path
                pipeline_options.ocr_options = ocr_options
                
            _log(f"OCR configured: {self.ocr_engine}")
        except Exception as e:
            _log(f"OCR configuration failed: {e}", level='warning')
            pipeline_options.do_ocr = False
    
    @property
    def is_available(self) -> bool:
        """Check if Docling is available and properly configured."""
        return self._docling_available
    
    @property
    def backend_name(self) -> str:
        """Return the name of the active backend."""
        return "docling" if self._docling_available else "legacy"
    
    @property
    def version(self) -> str:
        """Return Docling version."""
        return self._docling_version
    
    def get_status(self) -> Dict[str, Any]:
        """Get detailed status of the extractor."""
        return {
            'available': self._docling_available,
            'backend': self.backend_name,
            'version': self._docling_version,
            'artifacts_path': self.artifacts_path,
            'ocr_enabled': self.enable_ocr,
            'ocr_engine': self.ocr_engine if self.enable_ocr else None,
            'table_mode': self.table_mode,
            'offline_mode': True,  # Always true for this implementation
            'image_processing': False,  # Always disabled for memory optimization
        }
    
    def extract(self, filepath: str) -> DocumentExtractionResult:
        """
        Extract content from a document with maximum quality.
        
        Args:
            filepath: Path to the document file
            
        Returns:
            DocumentExtractionResult with all extracted content
        """
        start_time = time.time()
        
        filepath = str(Path(filepath).resolve())
        filename = os.path.basename(filepath)
        ext = Path(filepath).suffix.lower()
        
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {ext}. Supported: {self.SUPPORTED_FORMATS}")
        
        result = DocumentExtractionResult(
            filepath=filepath,
            filename=filename,
            format=ext.lstrip('.'),
            backend_used="unknown",
            offline_mode=True
        )
        
        try:
            if self._docling_available:
                result = self._extract_with_docling(filepath, result)
            elif self.fallback_to_legacy:
                result = self._extract_with_legacy(filepath, result)
            else:
                raise RuntimeError("Docling not available and fallback disabled")
                
        except Exception as e:
            _log(f"Extraction error: {e}", level='error')
            if self.fallback_to_legacy and result.backend_used != "legacy":
                _log("Falling back to legacy extraction...", level='warning')
                result.warnings.append(f"Docling failed: {e}. Using legacy extraction.")
                try:
                    result = self._extract_with_legacy(filepath, result)
                except Exception as e2:
                    result.warnings.append(f"Legacy extraction also failed: {e2}")
                    raise
            else:
                raise
        
        # Calculate final metrics
        result.extraction_time_ms = (time.time() - start_time) * 1000
        result.word_count = len(result.full_text.split())
        result.char_count = len(result.full_text)
        
        _log(f"Extracted {result.filename}: {result.word_count} words, "
             f"{len(result.tables)} tables, {len(result.paragraphs)} paragraphs, "
             f"{result.extraction_time_ms:.0f}ms (backend: {result.backend_used})")
        
        return result
    
    def _extract_with_docling(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract using Docling with full optimization."""
        result.backend_used = "docling"
        result.docling_version = self._docling_version
        result.models_used = ['layout', 'tableformer']
        if self.enable_ocr:
            result.models_used.append(f'ocr-{self.ocr_engine}')
        
        # Convert document
        conv_result = self._converter.convert(filepath)
        doc = conv_result.document
        
        # === Extract full text ===
        result.full_text = doc.export_to_text()
        
        # === Extract markdown (preserves structure) ===
        try:
            result.markdown = doc.export_to_markdown()
        except Exception as e:
            result.warnings.append(f"Markdown export failed: {e}")
            result.markdown = result.full_text
        
        # === Get page count ===
        try:
            if hasattr(conv_result, 'pages'):
                result.page_count = len(conv_result.pages)
            elif hasattr(doc, 'pages'):
                result.page_count = len(doc.pages)
            else:
                result.page_count = 1
        except:
            result.page_count = 1
        
        # === Extract tables with full detail ===
        table_id = 0
        try:
            for item in doc.iterate_items():
                item_obj = item[0] if isinstance(item, tuple) else item
                
                if self._is_table_item(item_obj):
                    table_id += 1
                    table = self._extract_table_data(item_obj, table_id)
                    if table:
                        result.tables.append(table)
        except Exception as e:
            result.warnings.append(f"Table extraction error: {e}")
        
        # === Extract sections/headings ===
        try:
            for item in doc.iterate_items():
                item_obj = item[0] if isinstance(item, tuple) else item
                section = self._extract_section_data(item_obj)
                if section:
                    result.sections.append(section)
        except Exception as e:
            result.warnings.append(f"Section extraction error: {e}")
        
        # === Extract all text paragraphs with location metadata ===
        para_num = 0
        try:
            for item in doc.iterate_items():
                item_obj = item[0] if isinstance(item, tuple) else item
                paragraphs = self._extract_paragraph_data(item_obj, para_num)
                for p in paragraphs:
                    para_num += 1
                    result.paragraphs.append(p)
        except Exception as e:
            result.warnings.append(f"Paragraph extraction error: {e}")
        
        # === Also extract text from tables for role extraction ===
        for table in result.tables:
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row):
                    if cell and cell.strip():
                        para_num += 1
                        result.paragraphs.append(ExtractedParagraph(
                            text=cell.strip(),
                            location=f"Table {table.table_id}, Row {row_idx+1}, Col {col_idx+1}",
                            page_number=table.page_number,
                            paragraph_type="table_cell",
                            confidence=table.confidence
                        ))
        
        return result
    
    def _is_table_item(self, item) -> bool:
        """Check if an item is a table."""
        if hasattr(item, 'label'):
            label = str(item.label).lower()
            return 'table' in label
        if hasattr(item, '__class__'):
            class_name = item.__class__.__name__.lower()
            return 'table' in class_name
        return False
    
    def _extract_table_data(self, item, table_id: int) -> Optional[ExtractedTable]:
        """Extract table data from a Docling table item."""
        try:
            headers = []
            rows = []
            confidence = 0.9  # Default high confidence for Docling tables
            
            # Try DataFrame export first (most reliable)
            if hasattr(item, 'export_to_dataframe'):
                try:
                    df = item.export_to_dataframe()
                    headers = [str(c) for c in df.columns.tolist()]
                    rows = [[str(cell) for cell in row] for row in df.values.tolist()]
                except:
                    pass
            
            # Fallback to raw data
            if not rows and hasattr(item, 'data'):
                data = item.data
                if isinstance(data, list) and data:
                    if isinstance(data[0], list):
                        headers = [str(c) for c in data[0]] if data else []
                        rows = [[str(cell) for cell in row] for row in data[1:]]
                    else:
                        rows = [[str(cell) for cell in data]]
            
            # Get metadata
            page_num = getattr(item, 'page_no', 1) or 1
            caption = getattr(item, 'caption', '') or ''
            if hasattr(item, 'confidence'):
                confidence = float(item.confidence)
            
            if rows or headers:
                return ExtractedTable(
                    table_id=table_id,
                    page_number=page_num,
                    headers=headers,
                    rows=rows,
                    caption=caption,
                    confidence=confidence
                )
        except Exception as e:
            _log(f"Table {table_id} extraction error: {e}", level='debug')
        
        return None
    
    def _extract_section_data(self, item) -> Optional[ExtractedSection]:
        """Extract section/heading data."""
        if not hasattr(item, 'label'):
            return None
        
        label = str(item.label).lower()
        if 'heading' not in label and 'title' not in label:
            return None
        
        # Determine heading level
        level = 1
        for i in range(1, 7):
            if f'heading_{i}' in label or f'heading{i}' in label or f'h{i}' in label:
                level = i
                break
        
        text = ""
        if hasattr(item, 'text'):
            text = item.text.strip() if item.text else ""
        elif hasattr(item, '__str__'):
            text = str(item).strip()
        
        if not text:
            return None
        
        page_num = getattr(item, 'page_no', 1) or 1
        
        return ExtractedSection(
            level=level,
            title=text,
            content="",
            page_number=page_num
        )
    
    def _extract_paragraph_data(self, item, current_num: int) -> List[ExtractedParagraph]:
        """Extract paragraph data with type classification."""
        paragraphs = []
        
        if not hasattr(item, 'text'):
            return paragraphs
        
        text = item.text.strip() if item.text else ""
        if not text:
            return paragraphs
        
        # Determine paragraph type
        para_type = "text"
        if hasattr(item, 'label'):
            label = str(item.label).lower()
            if 'heading' in label or 'title' in label:
                para_type = "heading"
            elif 'list' in label:
                para_type = "list_item"
            elif 'caption' in label:
                para_type = "caption"
        
        page_num = getattr(item, 'page_no', 1) or 1
        confidence = getattr(item, 'confidence', 1.0) or 1.0
        
        paragraphs.append(ExtractedParagraph(
            text=text,
            location=f"Page {page_num}, Block {current_num + 1}",
            page_number=page_num,
            paragraph_type=para_type,
            confidence=float(confidence)
        ))
        
        return paragraphs
    
    def _extract_with_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Fallback extraction using pdfplumber/python-docx."""
        result.backend_used = "legacy"
        result.docling_version = ""
        result.models_used = []
        ext = result.format
        
        if ext == 'pdf':
            result = self._extract_pdf_legacy(filepath, result)
        elif ext in ('docx', 'doc'):
            result = self._extract_docx_legacy(filepath, result)
        elif ext in ('pptx', 'ppt'):
            result = self._extract_pptx_legacy(filepath, result)
        elif ext in ('xlsx', 'xls'):
            result = self._extract_xlsx_legacy(filepath, result)
        elif ext in ('txt', 'md'):
            result = self._extract_txt(filepath, result)
        elif ext in ('html', 'htm'):
            result = self._extract_html_legacy(filepath, result)
        else:
            result.warnings.append(f"Legacy extraction not implemented for {ext}")
        
        return result
    
    def _extract_pdf_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract PDF using pdfplumber."""
        try:
            import pdfplumber
        except ImportError:
            try:
                import PyPDF2
                return self._extract_pdf_pypdf2(filepath, result)
            except ImportError:
                raise ImportError("pdfplumber or PyPDF2 required for PDF extraction")
        
        all_text = []
        para_num = 0
        
        with pdfplumber.open(filepath) as pdf:
            result.page_count = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                # Extract text
                text = page.extract_text() or ""
                if text:
                    all_text.append(text)
                    para_num += 1
                    result.paragraphs.append(ExtractedParagraph(
                        text=text,
                        location=f"Page {page_num}",
                        page_number=page_num,
                        paragraph_type="text"
                    ))
                
                # Extract tables
                tables = page.extract_tables() or []
                for table_idx, table_data in enumerate(tables):
                    if table_data and len(table_data) > 0:
                        headers = [str(cell or '') for cell in table_data[0]] if table_data else []
                        rows = [[str(cell or '') for cell in row] for row in table_data[1:]]
                        
                        result.tables.append(ExtractedTable(
                            table_id=len(result.tables) + 1,
                            page_number=page_num,
                            headers=headers,
                            rows=rows
                        ))
                        
                        # Add table cells to paragraphs
                        for row_idx, row in enumerate(rows):
                            for col_idx, cell in enumerate(row):
                                if cell.strip():
                                    para_num += 1
                                    result.paragraphs.append(ExtractedParagraph(
                                        text=cell.strip(),
                                        location=f"Page {page_num}, Table {table_idx+1}, Row {row_idx+1}",
                                        page_number=page_num,
                                        paragraph_type="table_cell"
                                    ))
        
        result.full_text = "\n\n".join(all_text)
        return result
    
    def _extract_pdf_pypdf2(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract PDF using PyPDF2 (no table support)."""
        import PyPDF2
        
        all_text = []
        para_num = 0
        
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            result.page_count = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text:
                    all_text.append(text)
                    para_num += 1
                    result.paragraphs.append(ExtractedParagraph(
                        text=text,
                        location=f"Page {page_num}",
                        page_number=page_num,
                        paragraph_type="text"
                    ))
        
        result.full_text = "\n\n".join(all_text)
        result.warnings.append("PyPDF2 fallback: table extraction unavailable")
        return result
    
    def _extract_docx_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract DOCX using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx required for DOCX extraction")
        
        doc = Document(filepath)
        all_text = []
        para_num = 0
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                para_num += 1
                all_text.append(text)
                
                # Determine type from style
                para_type = "text"
                if paragraph.style and paragraph.style.name:
                    style_name = paragraph.style.name.lower()
                    if 'heading' in style_name:
                        para_type = "heading"
                        # Extract section
                        level = 1
                        for i in range(1, 10):
                            if f'heading {i}' in style_name:
                                level = i
                                break
                        result.sections.append(ExtractedSection(
                            level=level,
                            title=text,
                            content="",
                            page_number=1
                        ))
                    elif 'list' in style_name:
                        para_type = "list_item"
                
                result.paragraphs.append(ExtractedParagraph(
                    text=text,
                    location=f"Paragraph {para_num}",
                    page_number=1,
                    paragraph_type=para_type
                ))
        
        # Extract tables
        for table_num, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows.append(row_data)
            
            if rows:
                result.tables.append(ExtractedTable(
                    table_id=table_num,
                    page_number=1,
                    headers=rows[0] if rows else [],
                    rows=rows[1:] if len(rows) > 1 else []
                ))
                
                # Add table cells to paragraphs
                for row_idx, row in enumerate(rows):
                    for col_idx, cell in enumerate(row):
                        if cell.strip():
                            para_num += 1
                            result.paragraphs.append(ExtractedParagraph(
                                text=cell.strip(),
                                location=f"Table {table_num}, Row {row_idx+1}, Col {col_idx+1}",
                                page_number=1,
                                paragraph_type="table_cell"
                            ))
        
        result.full_text = "\n\n".join(all_text)
        result.page_count = 1
        return result
    
    def _extract_pptx_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract PPTX using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            result.warnings.append("python-pptx not available for PPTX extraction")
            return result
        
        prs = Presentation(filepath)
        all_text = []
        para_num = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    para_num += 1
                    all_text.append(text)
                    result.paragraphs.append(ExtractedParagraph(
                        text=text,
                        location=f"Slide {slide_num}",
                        page_number=slide_num,
                        paragraph_type="text"
                    ))
        
        result.full_text = "\n\n".join(all_text)
        result.page_count = len(prs.slides)
        return result
    
    def _extract_xlsx_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract XLSX using openpyxl."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            result.warnings.append("openpyxl not available for XLSX extraction")
            return result
        
        wb = load_workbook(filepath, data_only=True)
        all_text = []
        para_num = 0
        
        for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
            sheet = wb[sheet_name]
            rows = []
            
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                if any(cell.strip() for cell in row_data):
                    rows.append(row_data)
            
            if rows:
                result.tables.append(ExtractedTable(
                    table_id=sheet_num,
                    page_number=sheet_num,
                    headers=rows[0] if rows else [],
                    rows=rows[1:] if len(rows) > 1 else [],
                    caption=sheet_name
                ))
                
                # Add cells to paragraphs
                for row_idx, row in enumerate(rows):
                    for col_idx, cell in enumerate(row):
                        if cell.strip():
                            para_num += 1
                            all_text.append(cell)
                            result.paragraphs.append(ExtractedParagraph(
                                text=cell.strip(),
                                location=f"Sheet '{sheet_name}', Row {row_idx+1}, Col {col_idx+1}",
                                page_number=sheet_num,
                                paragraph_type="table_cell"
                            ))
        
        result.full_text = "\n".join(all_text)
        result.page_count = len(wb.sheetnames)
        return result
    
    def _extract_html_legacy(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract HTML using BeautifulSoup or basic parsing."""
        try:
            from bs4 import BeautifulSoup
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
            
            # Remove script and style elements
            for element in soup(['script', 'style']):
                element.decompose()
            
            result.full_text = soup.get_text(separator='\n', strip=True)
            
            # Extract paragraphs
            para_num = 0
            for p in soup.find_all(['p', 'div', 'span', 'li']):
                text = p.get_text(strip=True)
                if text:
                    para_num += 1
                    result.paragraphs.append(ExtractedParagraph(
                        text=text,
                        location=f"Element {para_num}",
                        page_number=1,
                        paragraph_type="text"
                    ))
            
            # Extract headings
            for level in range(1, 7):
                for h in soup.find_all(f'h{level}'):
                    text = h.get_text(strip=True)
                    if text:
                        result.sections.append(ExtractedSection(
                            level=level,
                            title=text,
                            content="",
                            page_number=1
                        ))
            
        except ImportError:
            # Basic HTML extraction without BeautifulSoup
            import re
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read()
            
            # Strip tags
            text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
            
            result.full_text = text
            if text:
                result.paragraphs.append(ExtractedParagraph(
                    text=text,
                    location="Document",
                    page_number=1,
                    paragraph_type="text"
                ))
        
        result.page_count = 1
        return result
    
    def _extract_txt(self, filepath: str, result: DocumentExtractionResult) -> DocumentExtractionResult:
        """Extract plain text file."""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        
        result.full_text = text
        result.page_count = 1
        
        # Split into paragraphs
        paragraphs = text.split('\n\n')
        for para_num, para in enumerate(paragraphs, 1):
            para = para.strip()
            if para:
                result.paragraphs.append(ExtractedParagraph(
                    text=para,
                    location=f"Paragraph {para_num}",
                    page_number=1,
                    paragraph_type="text"
                ))
        
        return result


# ============================================================================
# DOCLING MANAGER - INSTALLATION & MODEL MANAGEMENT
# ============================================================================

class DoclingManager:
    """
    Utility class for managing Docling installation and models.
    
    Provides methods for:
    - Checking installation status
    - Downloading models for air-gapped use
    - Verifying offline configuration
    """
    
    # Windows-compatible path using os.path.join
    DEFAULT_MODELS_DIR = os.path.join(".cache", "docling", "models")
    
    REQUIRED_MODEL_DIRS = [
        'ds4sd--docling-models',  # Core layout/table models
    ]
    
    OPTIONAL_MODEL_DIRS = [
        'EasyOcr',  # EasyOCR models (if OCR enabled)
        'ds4sd--CodeFormulaV2',  # Code/formula detection
        'ds4sd--DocumentFigureClassifier',  # Figure classification
    ]
    
    @staticmethod
    def check_installation() -> Dict[str, Any]:
        """
        Check Docling installation status.
        
        Returns comprehensive status including:
        - Package installation
        - Model availability
        - Offline readiness
        """
        result = {
            'installed': False,
            'version': None,
            'pytorch_available': False,
            'pytorch_version': None,
            'models_path': None,
            'models_downloaded': False,
            'required_models': {},
            'optional_models': {},
            'offline_ready': False,
            'env_vars': {
                'DOCLING_ARTIFACTS_PATH': os.environ.get('DOCLING_ARTIFACTS_PATH'),
                'HF_HUB_OFFLINE': os.environ.get('HF_HUB_OFFLINE'),
                'TRANSFORMERS_OFFLINE': os.environ.get('TRANSFORMERS_OFFLINE'),
            }
        }
        
        # Check PyTorch
        try:
            import torch
            result['pytorch_available'] = True
            result['pytorch_version'] = torch.__version__
        except ImportError:
            pass
        
        # Check Docling
        try:
            import docling
            result['installed'] = True
            result['version'] = getattr(docling, '__version__', 'unknown')
        except ImportError:
            return result
        
        # Find models path
        models_path = (
            os.environ.get('DOCLING_ARTIFACTS_PATH') or
            os.environ.get('DOCLING_SERVE_ARTIFACTS_PATH')
        )
        
        if not models_path:
            # Check default location
            home = os.path.expanduser("~")
            default_path = os.path.join(home, DoclingManager.DEFAULT_MODELS_DIR)
            if os.path.exists(default_path):
                models_path = default_path
        
        if models_path and os.path.exists(models_path):
            result['models_path'] = models_path
            
            # Check required models
            for model_dir in DoclingManager.REQUIRED_MODEL_DIRS:
                path = os.path.join(models_path, model_dir)
                result['required_models'][model_dir] = os.path.exists(path)
            
            # Check optional models
            for model_dir in DoclingManager.OPTIONAL_MODEL_DIRS:
                path = os.path.join(models_path, model_dir)
                result['optional_models'][model_dir] = os.path.exists(path)
            
            result['models_downloaded'] = all(result['required_models'].values())
        
        # Determine offline readiness
        result['offline_ready'] = (
            result['installed'] and
            result['pytorch_available'] and
            result['models_downloaded']
        )
        
        return result
    
    @staticmethod
    def download_models(output_path: Optional[str] = None, include_ocr: bool = True) -> Dict[str, Any]:
        """
        Download Docling models for offline use.
        
        Args:
            output_path: Where to save models. If None, uses default cache.
            include_ocr: Whether to include OCR models (adds ~500MB)
            
        Returns:
            Dict with download status and paths
        """
        result = {
            'success': False,
            'output_path': output_path,
            'models_downloaded': [],
            'errors': []
        }
        
        try:
            # Try using docling-tools CLI approach
            from docling.utils.model_downloader import download_models as docling_download
            
            if output_path:
                os.makedirs(output_path, exist_ok=True)
                docling_download(output_path)
            else:
                docling_download()
            
            result['success'] = True
            result['models_downloaded'].append('core-models')
            _log(f"Models downloaded to: {output_path or 'default cache'}")
            
        except Exception as e:
            result['errors'].append(f"Model download failed: {e}")
            _log(f"Failed to download models: {e}", level='error')
        
        return result
    
    @staticmethod
    def get_package_requirements() -> Dict[str, List[str]]:
        """Return pip package requirements for Docling."""
        return {
            'core': [
                'docling>=2.70.0',
                'docling-core>=2.50.0',
                'docling-parse>=4.7.0',
                'docling-ibm-models>=3.9.0',
            ],
            'pytorch_cpu': [
                'torch',
                'torchvision',
            ],
            'ocr_easyocr': [
                'easyocr>=1.7.0',
            ],
            'ocr_tesseract': [
                'pytesseract>=0.3.10',
            ],
        }
    
    @staticmethod
    def estimate_disk_space() -> Dict[str, str]:
        """Return estimated disk space requirements."""
        return {
            'pytorch_cpu': '~800 MB',
            'docling_packages': '~700 MB',
            'core_models': '~1.2 GB',
            'ocr_models': '~500 MB (optional)',
            'total_minimum': '~2.7 GB',
            'total_with_ocr': '~3.2 GB'
        }
    
    @staticmethod
    def create_offline_config(models_path: str) -> str:
        """
        Generate environment configuration for offline operation.
        
        Args:
            models_path: Path where models are stored
            
        Returns:
            String content for a .env or batch file
        """
        config = f"""# Docling Air-Gap Configuration
# Generated by TechWriterReview DoclingManager

# Model artifacts location
DOCLING_ARTIFACTS_PATH={models_path}
DOCLING_SERVE_ARTIFACTS_PATH={models_path}

# Force offline mode - prevents ALL network access
HF_HUB_OFFLINE=1
TRANSFORMERS_OFFLINE=1
HF_DATASETS_OFFLINE=1

# Disable telemetry
HF_HUB_DISABLE_TELEMETRY=1
DO_NOT_TRACK=1
ANONYMIZED_TELEMETRY=false
"""
        return config


# ============================================================================
# CONVENIENCE FUNCTIONS
# ============================================================================

def create_extractor(**kwargs) -> DoclingExtractor:
    """Factory function to create a properly configured extractor."""
    return DoclingExtractor(**kwargs)


def extract_document(filepath: str, **kwargs) -> DocumentExtractionResult:
    """
    One-shot document extraction.
    
    Args:
        filepath: Path to document
        **kwargs: Passed to DoclingExtractor
        
    Returns:
        DocumentExtractionResult
    """
    extractor = DoclingExtractor(**kwargs)
    return extractor.extract(filepath)


def check_docling_status() -> Dict[str, Any]:
    """Check Docling installation and readiness."""
    return DoclingManager.check_installation()


# ============================================================================
# CLI / TEST
# ============================================================================

def main():
    """Command-line interface for testing."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Docling Document Extractor')
    parser.add_argument('filepath', nargs='?', help='Document to extract')
    parser.add_argument('--status', action='store_true', help='Check installation status')
    parser.add_argument('--ocr', action='store_true', help='Enable OCR')
    parser.add_argument('--table-mode', choices=['fast', 'accurate'], default='accurate')
    parser.add_argument('--artifacts', help='Path to model artifacts')
    
    args = parser.parse_args()
    
    if args.status or not args.filepath:
        print("\n" + "="*60)
        print("Docling Installation Status")
        print("="*60)
        status = DoclingManager.check_installation()
        print(json.dumps(status, indent=2))
        
        if not args.filepath:
            return
    
    if args.filepath:
        print(f"\n{'='*60}")
        print(f"Extracting: {args.filepath}")
        print("="*60 + "\n")
        
        extractor = DoclingExtractor(
            artifacts_path=args.artifacts,
            enable_ocr=args.ocr,
            table_mode=args.table_mode
        )
        
        result = extractor.extract(args.filepath)
        
        print(f"Backend: {result.backend_used}")
        print(f"Pages: {result.page_count}")
        print(f"Words: {result.word_count}")
        print(f"Tables: {len(result.tables)}")
        print(f"Sections: {len(result.sections)}")
        print(f"Paragraphs: {len(result.paragraphs)}")
        print(f"Time: {result.extraction_time_ms:.0f}ms")
        
        if result.warnings:
            print(f"\nWarnings:")
            for w in result.warnings:
                print(f"  - {w}")
        
        print(f"\n--- First 500 chars ---")
        print(result.full_text[:500])
        
        if result.tables:
            print(f"\n--- First table ---")
            t = result.tables[0]
            print(f"Headers: {t.headers}")
            print(f"Rows: {t.row_count}")


if __name__ == "__main__":
    main()

